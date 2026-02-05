#!/usr/bin/env python3
"""
Slide Layout Builders

High-quality, reusable slide creation functions.
Each function creates one slide type with consistent styling.
"""

import os
from typing import List, Optional, Literal, Dict, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from theme import PCPTheme, get_theme, apply_text_style, apply_title_style
from geometry import Box, fit_image, center_in_box, split_horizontal, check_text_overflow


def _add_title_textbox(
    slide,
    theme: PCPTheme,
    title: str,
    subtitle: Optional[str] = None
) -> None:
    """Add consistent title to any slide."""
    g = theme.geometry

    # Title textbox
    title_box = slide.shapes.add_textbox(
        g.margin_left,
        g.title_top,
        g.content_width,
        g.title_height
    )
    tf = title_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = title
    apply_title_style(p, theme)

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        apply_text_style(p2, theme, font_size=theme.typography.subtitle_size,
                        color=theme.colors.text_secondary)


def _get_body_box(theme: PCPTheme) -> Box:
    """Get the standard body area below title."""
    g = theme.geometry
    return Box(
        left=g.margin_left,
        top=g.body_top,
        width=g.content_width,
        height=Inches(g.slide_height.inches - g.body_top.inches - g.margin_bottom.inches)
    )


def add_title_slide(
    prs: Presentation,
    theme: PCPTheme,
    title: str,
    subtitle: Optional[str] = None,
    meta: Optional[str] = None
) -> None:
    """
    Add a title slide (opening slide).

    Args:
        prs: Presentation object
        theme: Theme configuration
        title: Main title text
        subtitle: Subtitle text
        meta: Additional metadata (date, author, etc.)
    """
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    g = theme.geometry

    # Center the title vertically
    title_top = Inches((g.slide_height.inches - 2) / 2)

    # Title
    title_box = slide.shapes.add_textbox(
        g.margin_left,
        title_top,
        g.content_width,
        Inches(1.2)
    )
    tf = title_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = title
    p.font.name = theme.typography.title_font
    p.font.size = theme.typography.title_size
    p.font.color.rgb = theme.colors.text_primary
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(
            g.margin_left,
            Inches(title_top.inches + 1.3),
            g.content_width,
            Inches(0.6)
        )
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.name = theme.typography.body_font
        p.font.size = theme.typography.subtitle_size
        p.font.color.rgb = theme.colors.text_secondary
        p.alignment = PP_ALIGN.CENTER

    # Meta (date, author)
    if meta:
        meta_box = slide.shapes.add_textbox(
            g.margin_left,
            Inches(g.slide_height.inches - 1),
            g.content_width,
            Inches(0.4)
        )
        tf = meta_box.text_frame
        p = tf.paragraphs[0]
        p.text = meta
        p.font.name = theme.typography.body_font
        p.font.size = theme.typography.caption_size
        p.font.color.rgb = theme.colors.text_muted
        p.alignment = PP_ALIGN.CENTER


def add_section_divider(
    prs: Presentation,
    theme: PCPTheme,
    title: str,
    kicker: Optional[str] = None
) -> None:
    """
    Add a section divider slide.

    Args:
        prs: Presentation object
        theme: Theme configuration
        title: Section title
        kicker: Optional small text above title
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    g = theme.geometry

    # Add accent bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(g.slide_height.inches / 2 - 0.1),
        Inches(0.15),
        Inches(0.8)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = theme.colors.accent
    bar.line.fill.background()

    # Kicker (small text above title)
    if kicker:
        kicker_box = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(g.slide_height.inches / 2 - 0.8),
            g.content_width,
            Inches(0.4)
        )
        tf = kicker_box.text_frame
        p = tf.paragraphs[0]
        p.text = kicker.upper()
        p.font.name = theme.typography.body_font
        p.font.size = theme.typography.caption_size
        p.font.color.rgb = theme.colors.accent
        p.font.bold = True

    # Section title
    title_box = slide.shapes.add_textbox(
        Inches(0.5),
        Inches(g.slide_height.inches / 2 - 0.3),
        g.content_width,
        Inches(1.2)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = theme.typography.title_font
    p.font.size = theme.typography.section_size
    p.font.color.rgb = theme.colors.text_primary
    p.font.bold = True


def add_bullets_slide(
    prs: Presentation,
    theme: PCPTheme,
    title: str,
    bullets: List[str],
    note: Optional[str] = None
) -> None:
    """
    Add a bullet points slide.

    Args:
        prs: Presentation object
        theme: Theme configuration
        title: Slide title (assertion)
        bullets: List of bullet points (max 5 recommended)
        note: Speaker notes
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    _add_title_textbox(slide, theme, title)
    body = _get_body_box(theme)

    # Bullets textbox
    bullet_box = slide.shapes.add_textbox(
        body.left,
        body.top,
        body.width,
        body.height
    )
    tf = bullet_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets[:theme.max_bullets]):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = f"  {bullet}"  # Indent
        p.font.name = theme.typography.body_font
        p.font.size = theme.typography.body_size
        p.font.color.rgb = theme.colors.text_primary
        p.space_before = theme.typography.paragraph_spacing_before
        p.space_after = theme.typography.paragraph_spacing_after

        # Add bullet character
        p.level = 0

    # Log warning if too many bullets
    if len(bullets) > theme.max_bullets:
        print(f"WARNING: Slide '{title[:30]}...' has {len(bullets)} bullets (max {theme.max_bullets})")

    # Speaker notes
    if note:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = note


def add_figure_slide(
    prs: Presentation,
    theme: PCPTheme,
    title: str,
    image_path: str,
    caption: Optional[str] = None,
    note: Optional[str] = None
) -> None:
    """
    Add a full-width figure slide.

    Args:
        prs: Presentation object
        theme: Theme configuration
        title: Slide title (assertion about the figure)
        image_path: Path to image file
        caption: Optional caption below image
        note: Speaker notes
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    _add_title_textbox(slide, theme, title)
    body = _get_body_box(theme)

    # Reserve space for caption if needed
    if caption:
        image_height = Inches(body.height.inches - 0.5)
    else:
        image_height = body.height

    # Check image exists
    if not os.path.exists(image_path):
        print(f"WARNING: Image not found: {image_path}")
        # Add placeholder
        placeholder = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            body.left,
            body.top,
            body.width,
            image_height
        )
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = theme.colors.background_alt
        return

    # Fit image maintaining aspect ratio
    img_width, img_height = fit_image(image_path, body.width, image_height)

    # Center image
    left, top = center_in_box(
        img_width, img_height,
        Box(body.left, body.top, body.width, image_height)
    )

    slide.shapes.add_picture(image_path, left, top, img_width, img_height)

    # Caption
    if caption:
        caption_box = slide.shapes.add_textbox(
            body.left,
            Inches(body.top.inches + image_height.inches + 0.1),
            body.width,
            Inches(0.4)
        )
        tf = caption_box.text_frame
        p = tf.paragraphs[0]
        p.text = caption
        p.font.name = theme.typography.body_font
        p.font.size = theme.typography.caption_size
        p.font.color.rgb = theme.colors.text_muted
        p.alignment = PP_ALIGN.CENTER

    # Speaker notes
    if note:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = note


def add_figure_with_text_slide(
    prs: Presentation,
    theme: PCPTheme,
    title: str,
    image_path: str,
    bullets: List[str],
    side: Literal["left", "right"] = "right",
    note: Optional[str] = None
) -> None:
    """
    Add a figure + bullet points slide (two-column).

    Args:
        prs: Presentation object
        theme: Theme configuration
        title: Slide title
        image_path: Path to image
        bullets: Bullet points
        side: Which side for the image ("left" or "right")
        note: Speaker notes
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    _add_title_textbox(slide, theme, title)
    body = _get_body_box(theme)

    # Split into two columns
    left_box, right_box = split_horizontal(body, ratio=0.5, gutter=theme.geometry.gutter)

    if side == "left":
        image_box, text_box = left_box, right_box
    else:
        text_box, image_box = left_box, right_box

    # Add image
    if os.path.exists(image_path):
        img_width, img_height = fit_image(image_path, image_box.width, image_box.height)
        left, top = center_in_box(img_width, img_height, image_box)
        slide.shapes.add_picture(image_path, left, top, img_width, img_height)
    else:
        print(f"WARNING: Image not found: {image_path}")

    # Add bullets
    bullet_box = slide.shapes.add_textbox(
        text_box.left,
        text_box.top,
        text_box.width,
        text_box.height
    )
    tf = bullet_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets[:theme.max_bullets]):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = f"  {bullet}"
        p.font.name = theme.typography.body_font
        p.font.size = theme.typography.body_size
        p.font.color.rgb = theme.colors.text_primary
        p.space_before = theme.typography.paragraph_spacing_before
        p.space_after = theme.typography.paragraph_spacing_after

    # Speaker notes
    if note:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = note


def add_comparison_slide(
    prs: Presentation,
    theme: PCPTheme,
    title: str,
    left: Dict[str, Any],
    right: Dict[str, Any],
    verdict: Optional[str] = None,
    note: Optional[str] = None
) -> None:
    """
    Add a comparison slide (A vs B, before/after).

    Args:
        prs: Presentation object
        theme: Theme configuration
        title: Slide title
        left: {"header": str, "points": [str]}
        right: {"header": str, "points": [str]}
        verdict: Optional conclusion line
        note: Speaker notes
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    _add_title_textbox(slide, theme, title)
    body = _get_body_box(theme)

    # Adjust for verdict if present
    if verdict:
        compare_height = Inches(body.height.inches - 0.6)
    else:
        compare_height = body.height

    compare_body = Box(body.left, body.top, body.width, compare_height)
    left_box, right_box = split_horizontal(compare_body, ratio=0.5, gutter=theme.geometry.gutter)

    # Add columns
    for box, data, color in [
        (left_box, left, theme.colors.text_secondary),
        (right_box, right, theme.colors.accent)
    ]:
        # Header
        header_box = slide.shapes.add_textbox(
            box.left,
            box.top,
            box.width,
            Inches(0.5)
        )
        tf = header_box.text_frame
        p = tf.paragraphs[0]
        p.text = data.get("header", "")
        p.font.name = theme.typography.body_font
        p.font.size = theme.typography.subtitle_size
        p.font.color.rgb = color
        p.font.bold = True

        # Points
        points_box = slide.shapes.add_textbox(
            box.left,
            Inches(box.top.inches + 0.6),
            box.width,
            Inches(box.height.inches - 0.6)
        )
        tf = points_box.text_frame
        tf.word_wrap = True

        for i, point in enumerate(data.get("points", [])):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.text = f"  {point}"
            p.font.name = theme.typography.body_font
            p.font.size = theme.typography.body_size
            p.font.color.rgb = theme.colors.text_primary

    # Verdict
    if verdict:
        verdict_box = slide.shapes.add_textbox(
            body.left,
            Inches(body.top.inches + compare_height.inches + 0.1),
            body.width,
            Inches(0.4)
        )
        tf = verdict_box.text_frame
        p = tf.paragraphs[0]
        p.text = verdict
        p.font.name = theme.typography.body_font
        p.font.size = theme.typography.body_size
        p.font.color.rgb = theme.colors.accent
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

    # Speaker notes
    if note:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = note


def add_key_number_slide(
    prs: Presentation,
    theme: PCPTheme,
    title: str,
    number: str,
    qualifier: Optional[str] = None,
    note: Optional[str] = None
) -> None:
    """
    Add a key number/metric slide.

    Args:
        prs: Presentation object
        theme: Theme configuration
        title: Slide title (assertion about the number)
        number: The key number/metric to display large
        qualifier: Context for the number
        note: Speaker notes
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    _add_title_textbox(slide, theme, title)
    body = _get_body_box(theme)

    # Big number in center
    number_box = slide.shapes.add_textbox(
        body.left,
        Inches(body.top.inches + (body.height.inches - 1.5) / 2),
        body.width,
        Inches(1.2)
    )
    tf = number_box.text_frame
    p = tf.paragraphs[0]
    p.text = number
    p.font.name = theme.typography.title_font
    p.font.size = theme.typography.key_number_size
    p.font.color.rgb = theme.colors.accent
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Qualifier below
    if qualifier:
        qual_box = slide.shapes.add_textbox(
            body.left,
            Inches(body.top.inches + (body.height.inches - 1.5) / 2 + 1.3),
            body.width,
            Inches(0.5)
        )
        tf = qual_box.text_frame
        p = tf.paragraphs[0]
        p.text = qualifier
        p.font.name = theme.typography.body_font
        p.font.size = theme.typography.subtitle_size
        p.font.color.rgb = theme.colors.text_secondary
        p.alignment = PP_ALIGN.CENTER

    # Speaker notes
    if note:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = note


def add_process_slide(
    prs: Presentation,
    theme: PCPTheme,
    title: str,
    steps: List[str],
    note: Optional[str] = None
) -> None:
    """
    Add a process/flow slide.

    Args:
        prs: Presentation object
        theme: Theme configuration
        title: Slide title
        steps: List of process steps
        note: Speaker notes
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    _add_title_textbox(slide, theme, title)
    body = _get_body_box(theme)

    num_steps = len(steps)
    if num_steps == 0:
        return

    # Calculate step box sizes
    step_width = (body.width.inches - (num_steps - 1) * 0.3) / num_steps
    arrow_width = 0.3

    for i, step in enumerate(steps):
        # Step box position
        step_left = body.left.inches + i * (step_width + arrow_width)

        # Step number circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(step_left + step_width / 2 - 0.25),
            body.top,
            Inches(0.5),
            Inches(0.5)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = theme.colors.accent
        circle.line.fill.background()

        # Number in circle
        num_box = slide.shapes.add_textbox(
            Inches(step_left + step_width / 2 - 0.25),
            body.top,
            Inches(0.5),
            Inches(0.5)
        )
        tf = num_box.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        p.text = str(i + 1)
        p.font.name = theme.typography.body_font
        p.font.size = Pt(16)
        p.font.color.rgb = theme.colors.text_inverse
        p.font.bold = True

        # Step text
        text_box = slide.shapes.add_textbox(
            Inches(step_left),
            Inches(body.top.inches + 0.7),
            Inches(step_width),
            Inches(body.height.inches - 0.7)
        )
        tf = text_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = step
        p.font.name = theme.typography.body_font
        p.font.size = theme.typography.body_size
        p.font.color.rgb = theme.colors.text_primary
        p.alignment = PP_ALIGN.CENTER

        # Arrow (except after last step)
        if i < num_steps - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                Inches(step_left + step_width + 0.05),
                Inches(body.top.inches + 0.1),
                Inches(0.2),
                Inches(0.3)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = theme.colors.text_muted
            arrow.line.fill.background()

    # Speaker notes
    if note:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = note


def add_table_slide(
    prs: Presentation,
    theme: PCPTheme,
    title: str,
    data: List[List[str]],
    headers: Optional[List[str]] = None,
    note: Optional[str] = None
) -> None:
    """
    Add a table slide.

    Args:
        prs: Presentation object
        theme: Theme configuration
        title: Slide title
        data: Table data as list of rows
        headers: Optional header row
        note: Speaker notes
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    _add_title_textbox(slide, theme, title)
    body = _get_body_box(theme)

    if not data:
        return

    # Calculate table dimensions
    rows = len(data) + (1 if headers else 0)
    cols = len(data[0]) if data else 0

    if cols == 0:
        return

    # Add table
    table = slide.shapes.add_table(
        rows, cols,
        body.left,
        body.top,
        body.width,
        Inches(min(body.height.inches, rows * 0.5))
    ).table

    # Style header row
    if headers:
        for j, header in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = theme.colors.primary

            p = cell.text_frame.paragraphs[0]
            p.font.name = theme.typography.body_font
            p.font.size = theme.typography.body_size
            p.font.color.rgb = theme.colors.text_inverse
            p.font.bold = True

    # Fill data rows
    start_row = 1 if headers else 0
    for i, row_data in enumerate(data):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(start_row + i, j)
            cell.text = str(cell_text)

            p = cell.text_frame.paragraphs[0]
            p.font.name = theme.typography.body_font
            p.font.size = theme.typography.body_size
            p.font.color.rgb = theme.colors.text_primary

    # Speaker notes
    if note:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = note


# Layout registry for build_deck.py
LAYOUT_FUNCTIONS = {
    "title": add_title_slide,
    "section": add_section_divider,
    "bullets": add_bullets_slide,
    "figure": add_figure_slide,
    "figure_with_text": add_figure_with_text_slide,
    "comparison": add_comparison_slide,
    "key_number": add_key_number_slide,
    "process": add_process_slide,
    "table": add_table_slide,
}


if __name__ == "__main__":
    print("Available slide layouts:")
    for name in LAYOUT_FUNCTIONS:
        print(f"  - {name}")
